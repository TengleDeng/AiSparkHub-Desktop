#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
文件查看器组件
用于在标签页中显示HTML、Markdown、Word、Excel、PowerPoint、PDF和文本文件内容
"""

import os
import markdown
import io
import tempfile
import qtawesome as qta
from PyQt6.QtCore import Qt, QUrl, QSize, pyqtSignal, pyqtSlot
from PyQt6.QtWidgets import (
    QWidget, QTextEdit, QVBoxLayout, QToolBar, QApplication,
    QPushButton, QLabel
)
from PyQt6.QtWebEngineWidgets import QWebEngineView
from PyQt6.QtGui import QFont, QIcon, QColor, QTextCharFormat, QTextCursor, QAction

# 导入文档处理库
try:
    import docx
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

try:
    import openpyxl
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False

try:
    import fitz  # PyMuPDF
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

class FileViewer(QWidget):
    """文件查看器组件，用于显示单个文件内容"""
    
    # 添加信号，用于将文件内容复制到提示词
    file_content_to_prompt = pyqtSignal(str)
    
    def __init__(self, theme_manager=None):
        super().__init__()
        
        # 获取或设置主题管理器
        self.theme_manager = theme_manager or QApplication.instance().theme_manager
        
        # 创建布局
        self.layout = QVBoxLayout(self)
        self.layout.setContentsMargins(0, 0, 0, 0)
        self.layout.setSpacing(0) # No spacing between toolbar and viewer
        
        # 创建工具栏
        self.toolbar = QToolBar()
        self.toolbar.setIconSize(QSize(16, 16))
        
        # 添加工具栏
        self.layout.addWidget(self.toolbar)
        
        # 初始化成员变量
        self.file_path = None
        self.file_type = None
        self.viewer = None
        self.copy_to_prompt_action = None # Store the action
        
        # 连接主题变化信号
        if self.theme_manager:
            self.theme_manager.theme_changed.connect(self._update_theme)
            self._update_theme() # Apply initial theme (styles toolbar)
    
    def open_file(self, file_path, file_type=None):
        """打开文件并显示内容
        
        Args:
            file_path (str): 文件路径
            file_type (str, optional): 文件类型。默认为None，会根据扩展名自动判断。
        """
        if not os.path.exists(file_path):
            print(f"File not found: {file_path}")
            return
        
        # 保存文件路径和类型
        self.file_path = file_path
        
        # 如果没有指定文件类型，根据扩展名判断
        if file_type is None:
            ext = os.path.splitext(file_path)[1].lower()
            if ext in ['.md', '.markdown']:
                file_type = 'markdown'
            elif ext in ['.html', '.htm']:
                file_type = 'html'
            elif ext in ['.docx', '.doc']:
                file_type = 'docx'
            elif ext in ['.pptx', '.ppt']:
                file_type = 'powerpoint'
            elif ext in ['.xlsx', '.xls']:
                file_type = 'excel'
            elif ext in ['.pdf']:
                file_type = 'pdf'
            else:
                file_type = 'text'
        
        self.file_type = file_type # Store determined file type
        
        # 清空工具栏 (先保留按钮实例)
        self.toolbar.clear()
        
        # 添加复制到提示词按钮 (或重新添加已存在的按钮)
        if self.copy_to_prompt_action is None:
            self.copy_to_prompt_action = QAction(qta.icon('fa5s.copy'), "复制到提示词", self)
            self.copy_to_prompt_action.triggered.connect(self._copy_to_prompt)
        self.toolbar.addAction(self.copy_to_prompt_action)
        
        # 如果已有查看器，移除它
        if self.viewer is not None:
            self.layout.removeWidget(self.viewer)
            self.viewer.deleteLater()
            self.viewer = None
        
        # 调用 _update_theme 来创建/更新查看器并应用主题
        self._update_theme()
    
    def _update_theme(self):
        """更新组件的主题样式和内容"""
        if not self.theme_manager:
            return
        
        theme_colors = self.theme_manager.get_current_theme_colors()
        
        # --- 更新工具栏样式 ---
        toolbar_bg = theme_colors.get('secondary_bg', '#3B4252')
        toolbar_hover_bg = theme_colors.get('tertiary_bg', '#4C566A')
        icon_color = theme_colors.get('foreground', '#D8DEE9')
        
        self.toolbar.setStyleSheet(f"""
            QToolBar {{
                background-color: {toolbar_bg};
                border: none;
                spacing: 5px;
                padding: 5px;
            }}
            QToolButton {{
                background-color: transparent;
                border: none;
                border-radius: 4px;
                padding: 5px;
                color: {icon_color}; /* Set icon/text color */
            }}
            QToolButton:hover {{
                background-color: {toolbar_hover_bg};
            }}
        """)
        
        # --- 更新复制按钮图标颜色 ---
        if self.copy_to_prompt_action:
            self.copy_to_prompt_action.setIcon(qta.icon('fa5s.copy', color=icon_color))
        
        # --- 更新查看器样式和内容 (仅当文件已加载时) ---
        if not self.file_path:
            return
        
        # 确定查看器类型
        if self.file_type in ['text']:
            viewer_type = QTextEdit
        elif self.file_type in ['html', 'markdown', 'docx', 'powerpoint', 'excel', 'pdf']:
            viewer_type = QWebEngineView
        else:
            # 如果类型未知或不支持，也使用文本编辑器显示错误或原始文本
            viewer_type = QTextEdit
            self.file_type = 'text' # Treat as text
        
        # 创建或重用查看器
        if self.viewer is None or not isinstance(self.viewer, viewer_type):
            if self.viewer is not None:
                old_viewer = self.viewer
                self.layout.removeWidget(old_viewer)
                old_viewer.deleteLater()
                self.viewer = None # Ensure it's None before creating new
        
        self.viewer = viewer_type()
        if isinstance(self.viewer, QTextEdit):
            self.viewer.setReadOnly(True)
            font = QFont("Consolas", 'Courier New')
            font.setStyleHint(QFont.StyleHint.Monospace)
            self.viewer.setFont(font)
        self.layout.addWidget(self.viewer)
        
        # --- 加载内容并应用样式 ---
        try:
            if isinstance(self.viewer, QTextEdit):
                text_bg = theme_colors.get('background', '#2E3440')
                text_fg = theme_colors.get('foreground', '#D8DEE9')
                self.viewer.setStyleSheet(f"""
                    QTextEdit {{
                        background-color: {text_bg};
                        color: {text_fg};
                        border: none;
                        padding: 8px;
                    }}
                """)
                # 读取文本内容
                content = self._get_text_content(self.file_path)
                self.viewer.setPlainText(content)
            
            elif isinstance(self.viewer, QWebEngineView):
                base_html_content = ""
                if self.file_type == 'html':
                    # For raw HTML, we could try injecting CSS, but for now just load it.
                    # Theming might not work perfectly for external HTML.
                    self.viewer.setUrl(QUrl.fromLocalFile(self.file_path))
                    return # Skip styled HTML generation for raw HTML
                elif self.file_type == 'markdown':
                    base_html_content = self._md_to_html(self.file_path)
                elif self.file_type == 'docx' and DOCX_SUPPORT:
                    base_html_content = self._docx_to_html(self.file_path)
                elif self.file_type == 'powerpoint' and PPTX_SUPPORT:
                    base_html_content = self._pptx_to_html(self.file_path)
                elif self.file_type == 'excel' and EXCEL_SUPPORT:
                    base_html_content = self._excel_to_html(self.file_path)
                elif self.file_type == 'pdf' and PDF_SUPPORT:
                    base_html_content = self._pdf_to_html(self.file_path)
                else:
                    base_html_content = f"<p>不支持的文件类型: {self.file_type}</p>"
                
                # Create styled HTML and set it
                styled_html = self._create_styled_html(base_html_content)
                # Use setHtml with a base URL for relative paths if needed
                self.viewer.setHtml(styled_html, QUrl.fromLocalFile(os.path.dirname(self.file_path)))
        
        except Exception as e:
            error_message = f"无法加载或应用主题: {e}"
            if isinstance(self.viewer, QTextEdit):
                self.viewer.setPlainText(error_message)
            elif isinstance(self.viewer, QWebEngineView):
                # Display error within the web view itself
                error_html = self._create_styled_html(f"<p style='color: red;'>{error_message}</p>")
                self.viewer.setHtml(error_html)
            print(f"Error updating file viewer theme/content: {e}")
    
    def _get_text_content(self, file_path):
        """安全地读取文本文件内容"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            try:
                # Fallback to system default encoding (less reliable)
                with open(file_path, 'r') as f:
                    return f.read()
            except Exception as decode_err:
                return f"无法解码文件内容 ({decode_err}): {file_path}"
        except Exception as read_err:
            return f"读取文件时出错 ({read_err}): {file_path}"
    
    def _md_to_html(self, file_path):
        """将Markdown文件转换为基础HTML"""
        try:
            md_content = self._get_text_content(file_path)
            if md_content.startswith("无法") or md_content.startswith("读取文件时出错"):
                return f"<p>{md_content}</p>"
            html_content = markdown.markdown(
                md_content,
                extensions=['tables', 'fenced_code', 'codehilite'] # Add more extensions if needed
            )
            return html_content
        except Exception as e:
            return f"<p>Markdown转换失败: {e}</p>"
    
    def _create_styled_html(self, content):
        """创建带有样式的HTML内容
        
        Args:
            content (str): HTML内容
            
        Returns:
            str: 带样式的HTML
        """
        if not self.theme_manager:
            # Fallback if theme manager isn't available
            # (You might want to define default dark/light colors here)
            bg_color = "#2E3440"
            fg_color = "#D8DEE9"
            secondary_bg = "#3B4252"
            accent_color = "#88C0D0"
        else:
            # 获取当前主题颜色
            theme_colors = self.theme_manager.get_current_theme_colors()
            bg_color = theme_colors.get('background', "#2E3440")
            fg_color = theme_colors.get('foreground', "#D8DEE9")
            secondary_bg = theme_colors.get('secondary_bg', "#3B4252")
            accent_color = theme_colors.get('accent', "#88C0D0")
            tertiary_bg = theme_colors.get('tertiary_bg', "#4C566A") # Used for borders etc.
        
        return f"""
        <html>
        <head>
            <style>
                body {{
                    font-family: 'Segoe UI', Arial, sans-serif;
                    line-height: 1.6;
                    color: {fg_color};
                    background-color: {bg_color};
                    padding: 20px;
                    max-width: 900px;
                    margin: 0 auto;
                }}
                h1, h2, h3, h4, h5, h6 {{
                    color: {fg_color}; /* Use main foreground color */
                    margin-top: 24px;
                    margin-bottom: 16px;
                }}
                h1 {{ font-size: 2em; border-bottom: 1px solid {tertiary_bg}; padding-bottom: 0.3em; }}
                h2 {{ font-size: 1.5em; border-bottom: 1px solid {tertiary_bg}; padding-bottom: 0.3em; }}
                a {{ color: {accent_color}; text-decoration: none; }}
                a:hover {{ text-decoration: underline; }}
                code {{
                    font-family: 'Consolas', 'Monaco', monospace;
                    background-color: {secondary_bg};
                    padding: 0.2em 0.4em;
                    border-radius: 3px;
                }}
                pre {{
                    background-color: {secondary_bg};
                    padding: 16px;
                    border-radius: 4px;
                    overflow: auto;
                }}
                pre code {{
                    background-color: transparent;
                    padding: 0;
                }}
                table {{
                    border-collapse: collapse;
                    width: 100%;
                    margin: 16px 0;
                }}
                table, th, td {{
                    border: 1px solid {tertiary_bg};
                }}
                th, td {{
                    padding: 8px 12px;
                    text-align: left;
                }}
                th {{
                    background-color: {secondary_bg};
                }}
                blockquote {{
                    border-left: 4px solid {tertiary_bg};
                    margin: 16px 0;
                    padding: 0 16px;
                    color: {fg_color};
                }}
                img {{
                    max-width: 100%;
                    height: auto;
                }}
                hr {{
                    border: none;
                    height: 1px;
                    background-color: {tertiary_bg};
                    margin: 24px 0;
                }}
            </style>
        </head>
        <body>
            {content}
        </body>
        </html>
        """
    
    def _docx_to_html(self, file_path):
        """将Word文档转换为HTML
        
        Args:
            file_path (str): Word文档路径
            
        Returns:
            str: HTML内容
        """
        if not DOCX_SUPPORT:
            return "<p>未安装python-docx库，无法查看Word文档</p>"
            
        doc = docx.Document(file_path)
        html = []
        
        # 处理文档内容
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
                
            # 处理不同样式的段落
            if para.style.name.startswith('Heading'):
                level = para.style.name.replace('Heading', '')
                try:
                    level_num = int(level.strip())
                    if 1 <= level_num <= 6:
                        html.append(f"<h{level_num}>{para.text}</h{level_num}>")
                        continue
                except ValueError:
                    pass
            
            # 默认作为普通段落处理
            html.append(f"<p>{para.text}</p>")
        
        # 处理表格
        for table in doc.tables:
            html_table = ["<table>"]
            
            # 处理表格行
            for i, row in enumerate(table.rows):
                html_table.append("<tr>")
                
                # 处理单元格
                for cell in row.cells:
                    # 第一行通常是表头
                    if i == 0:
                        html_table.append(f"<th>{cell.text}</th>")
                    else:
                        html_table.append(f"<td>{cell.text}</td>")
                
                html_table.append("</tr>")
            
            html_table.append("</table>")
            html.append("".join(html_table))
        
        return "".join(html)
    
    def _pptx_to_html(self, file_path):
        """将PowerPoint转换为HTML
        
        Args:
            file_path (str): PowerPoint文件路径
            
        Returns:
            str: HTML内容
        """
        if not PPTX_SUPPORT:
            return "<p>未安装python-pptx库，无法查看PowerPoint文件</p>"
            
        prs = Presentation(file_path)
        html = []
        
        # 处理每张幻灯片
        for i, slide in enumerate(prs.slides):
            html.append(f"<div class='slide'><h2>幻灯片 {i+1}</h2>")
            
            # 处理幻灯片标题
            if slide.shapes.title:
                html.append(f"<h3>{slide.shapes.title.text}</h3>")
            
            # 处理幻灯片内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip() and shape != slide.shapes.title:
                    html.append(f"<p>{shape.text}</p>")
            
            html.append("</div><hr>")
        
        return "".join(html)
    
    def _excel_to_html(self, file_path):
        """将Excel文件转换为HTML表格
        
        Args:
            file_path (str): Excel文件路径
            
        Returns:
            str: HTML内容
        """
        if not EXCEL_SUPPORT:
            return "<p>未安装openpyxl库，无法查看Excel文件</p>"
            
        wb = openpyxl.load_workbook(file_path, data_only=True)
        html = []
        
        # 处理每个工作表
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            html.append(f"<h2>{sheet_name}</h2>")
            html.append("<table>")
            
            # 获取使用的单元格范围
            min_row, min_col = 1, 1
            max_row = max(1, sheet.max_row)
            max_col = max(1, sheet.max_column)
            
            # 限制显示的最大行数，避免过大的表格导致性能问题
            max_display_rows = min(200, max_row)
            
            # 处理表格内容
            for row in range(min_row, min_row + max_display_rows):
                html.append("<tr>")
                
                for col in range(min_col, max_col + 1):
                    cell = sheet.cell(row=row, column=col)
                    cell_value = str(cell.value) if cell.value is not None else ""
                    
                    # 第一行通常是表头
                    if row == min_row:
                        html.append(f"<th>{cell_value}</th>")
                    else:
                        html.append(f"<td>{cell_value}</td>")
                
                html.append("</tr>")
            
            # 如果有更多行，显示提示信息
            if max_row > max_display_rows:
                html.append(
                    f"<tr><td colspan='{max_col}' style='text-align:center;'>...</td></tr>"
                    f"<tr><td colspan='{max_col}' style='text-align:center;'>（还有 {max_row - max_display_rows} 行未显示）</td></tr>"
                )
            
            html.append("</table>")
            html.append("<hr>")
        
        return "".join(html)
    
    def _pdf_to_html(self, file_path):
        """将PDF文件转换为HTML
        
        Args:
            file_path (str): PDF文件路径
            
        Returns:
            str: HTML内容
        """
        if not PDF_SUPPORT:
            return "<p>未安装PyMuPDF库，无法查看PDF文件</p>"
            
        # 打开PDF文件
        pdf_document = fitz.open(file_path)
        html = []
        
        # 处理每一页
        for page_num, page in enumerate(pdf_document):
            html.append(f"<div class='pdf-page'><h2>第 {page_num + 1} 页</h2>")
            
            # 获取页面文本
            text = page.get_text()
            # 替换换行符为<br>标签
            text = text.replace("\n", "<br>")
            
            html.append(f"<div class='pdf-content'>{text}</div>")
            html.append("</div><hr>")
        
        return "".join(html)
    
    def _copy_to_prompt(self):
        """复制文件内容到提示词输入框"""
        if not self.file_path:
            return
            
        try:
            with open(self.file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                # 发出信号，将内容传递给提示词输入框
                self.file_content_to_prompt.emit(content)
        except UnicodeDecodeError:
            try:
                # 尝试使用系统默认编码
                with open(self.file_path, 'r') as f:
                    content = f.read()
                    self.file_content_to_prompt.emit(content)
            except Exception as e:
                print(f"读取文件失败: {e}")
        except Exception as e:
            print(f"读取文件失败: {e}") 