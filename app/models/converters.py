#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
文件格式转换器 - 将不同格式的文件内容提取为纯文本并转换为Markdown格式
"""

import os
import re
import html
from abc import ABC, abstractmethod
from bs4 import BeautifulSoup
try:
    import markdown
    from markdown.extensions.tables import TableExtension
    from markdown.extensions.fenced_code import FencedCodeExtension
except ImportError:
    markdown = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None


class FormatConverter(ABC):
    """格式转换器基类"""
    
    @abstractmethod
    def extract_content(self, file_path):
        """
        从文件中提取内容
        
        Args:
            file_path (str): 文件路径
            
        Returns:
            tuple: (content, title) - 提取的内容和标题
        """
        pass
    
    @abstractmethod
    def convert_to_markdown(self, content):
        """
        将内容转换为Markdown格式
        
        Args:
            content (str): 原始内容
            
        Returns:
            str: Markdown格式的内容
        """
        pass


class MarkdownConverter(FormatConverter):
    """Markdown格式转换器"""
    
    def extract_content(self, file_path):
        """
        从Markdown文件中提取内容
        
        Args:
            file_path (str): 文件路径
            
        Returns:
            tuple: (content, title) - 提取的内容和标题
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 提取标题
            title = self._extract_title(content)
            if not title:
                # 如果没有提取到标题，使用文件名作为标题
                title = os.path.splitext(os.path.basename(file_path))[0]
                
            return content, title
        except Exception as e:
            print(f"提取Markdown内容出错: {e}")
            return "", os.path.splitext(os.path.basename(file_path))[0]
    
    def convert_to_markdown(self, content):
        """
        将内容转换为Markdown格式（Markdown本身无需转换）
        
        Args:
            content (str): 原始内容
            
        Returns:
            str: Markdown格式的内容
        """
        return content
    
    def _extract_title(self, content):
        """
        从Markdown内容中提取标题
        
        Args:
            content (str): Markdown内容
            
        Returns:
            str: 提取的标题，如果没有则返回None
        """
        # 尝试从# 标题行中提取
        title_match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
        if title_match:
            return title_match.group(1).strip()
        
        # 如果没有# 标题行，尝试从内容的第一行提取
        lines = content.split('\n')
        if lines and lines[0].strip():
            return lines[0].strip()
        
        return None


class HtmlConverter(FormatConverter):
    """HTML格式转换器"""
    
    def extract_content(self, file_path):
        """
        从HTML文件中提取内容
        
        Args:
            file_path (str): 文件路径
            
        Returns:
            tuple: (content, title) - 提取的内容和标题
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 提取标题
            title_tag = soup.find('title')
            title = title_tag.get_text().strip() if title_tag else None
            
            if not title:
                # 尝试从h1标签提取标题
                h1_tag = soup.find('h1')
                title = h1_tag.get_text().strip() if h1_tag else None
                
            if not title:
                # 如果仍然没有标题，使用文件名
                title = os.path.splitext(os.path.basename(file_path))[0]
            
            # 移除script和style标签
            for script in soup(["script", "style"]):
                script.extract()
            
            # 获取可见文本
            text = soup.get_text()
            
            # 格式化文本，移除多余的空白
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            # 从body区域提取原始HTML，用于更好的Markdown转换
            body = soup.find('body')
            html_for_conversion = str(body) if body else html_content
            
            return html_for_conversion, title
        except Exception as e:
            print(f"提取HTML内容出错: {e}")
            return "", os.path.splitext(os.path.basename(file_path))[0]
    
    def convert_to_markdown(self, content):
        """
        将HTML内容转换为Markdown格式
        
        Args:
            content (str): HTML内容
            
        Returns:
            str: Markdown格式的内容
        """
        try:
            # 使用html2text库转换HTML到Markdown
            import html2text
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = False
            h.ignore_tables = False
            h.ignore_emphasis = False
            
            # 将HTML转换为Markdown
            markdown_content = h.handle(content)
            return markdown_content
        except ImportError:
            # 如果html2text不可用，使用简单的BeautifulSoup转换
            try:
                soup = BeautifulSoup(content, 'html.parser')
                # 移除script和style标签
                for script in soup(["script", "style"]):
                    script.extract()
                
                # 获取文本
                text = soup.get_text()
                # 格式化文本
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = '\n'.join(chunk for chunk in chunks if chunk)
                
                return text
            except Exception as e:
                print(f"HTML到Markdown转换出错: {e}")
                return content


class TextConverter(FormatConverter):
    """纯文本格式转换器"""
    
    def extract_content(self, file_path):
        """
        从纯文本文件中提取内容
        
        Args:
            file_path (str): 文件路径
            
        Returns:
            tuple: (content, title) - 提取的内容和标题
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 使用文件名作为标题
            title = os.path.splitext(os.path.basename(file_path))[0]
            
            # 尝试从第一行提取标题
            lines = content.split('\n')
            if lines and lines[0].strip():
                title = lines[0].strip()
                
            return content, title
        except Exception as e:
            print(f"提取文本内容出错: {e}")
            return "", os.path.splitext(os.path.basename(file_path))[0]
    
    def convert_to_markdown(self, content):
        """
        将纯文本内容转换为Markdown格式
        
        Args:
            content (str): 纯文本内容
            
        Returns:
            str: Markdown格式的内容
        """
        # 纯文本基本可以兼容Markdown格式，只需要处理一些特殊字符
        # 转义Markdown特殊字符
        text = content
        for char in ['\\', '`', '*', '_', '{', '}', '[', ']', '(', ')', '#', '+', '-', '.', '!']:
            text = text.replace(char, '\\' + char)
        
        # 将换行符保留为Markdown硬换行
        text = text.replace('\n', '  \n')
        
        return text


class PdfConverter(FormatConverter):
    """PDF格式转换器"""
    
    def extract_content(self, file_path):
        """
        从PDF文件中提取内容
        
        Args:
            file_path (str): 文件路径
            
        Returns:
            tuple: (content, title) - 提取的内容和标题
        """
        if fitz is None and PyPDF2 is None:
            print("警告: PyMuPDF和PyPDF2库均未安装，无法提取PDF内容")
            return "", os.path.splitext(os.path.basename(file_path))[0]
        
        try:
            title = os.path.splitext(os.path.basename(file_path))[0]
            content = ""
            
            # 优先使用PyMuPDF (fitz)
            if fitz is not None:
                # 使用PyMuPDF提取文本
                doc = fitz.open(file_path)
                
                # 尝试从PDF元数据中提取标题
                if doc.metadata and doc.metadata.get('title'):
                    pdf_title = doc.metadata.get('title')
                    if pdf_title and pdf_title.strip():
                        title = pdf_title.strip()
                
                # 提取所有页面的文本内容
                page_texts = []
                toc = doc.get_toc()  # 获取目录结构，如果有的话
                
                # 如果有目录，使用它来构建更好的章节结构
                if toc:
                    # 构建目录结构
                    toc_dict = {}
                    for level, title, page in toc:
                        if page < len(doc) and page >= 0:
                            if page not in toc_dict:
                                toc_dict[page] = []
                            toc_dict[page].append((level, title))
                
                # 遍历所有页面
                for page_num, page in enumerate(doc):
                    # 如果这个页面有目录条目，先添加它们
                    if toc and page_num in toc_dict:
                        for level, heading in toc_dict[page_num]:
                            # 根据级别添加适当数量的#作为Markdown标题
                            page_texts.append(f"{'#' * min(level, 6)} {heading}\n\n")
                    
                    # 获取页面文本
                    text = page.get_text("text")
                    if text:
                        page_texts.append(text)
                
                content = "\n\n".join(page_texts)
                
                # 尝试识别段落和格式
                content = self._improve_text_structure(content)
                
                doc.close()
            else:
                # 回退到PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    
                    # 尝试从PDF元数据中提取标题
                    if reader.metadata and '/Title' in reader.metadata:
                        pdf_title = reader.metadata['/Title']
                        if pdf_title:
                            title = pdf_title
                    
                    # 提取所有页面的文本内容
                    for page_num in range(len(reader.pages)):
                        page = reader.pages[page_num]
                        content += page.extract_text() + "\n\n"
            
            return content, title
        except Exception as e:
            print(f"提取PDF内容出错: {e}")
            return "", os.path.splitext(os.path.basename(file_path))[0]
    
    def _improve_text_structure(self, text):
        """
        改进从PDF提取的文本结构，尝试识别标题、段落等
        
        Args:
            text (str): 原始提取的文本
            
        Returns:
            str: 改进后的文本结构
        """
        lines = text.split('\n')
        improved_lines = []
        current_paragraph = []
        
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                # 处理空行：结束当前段落
                if current_paragraph:
                    improved_lines.append(' '.join(current_paragraph))
                    current_paragraph = []
                improved_lines.append('')
                continue
            
            # 检测可能的标题
            if len(line) < 100 and line.strip() and (i == 0 or not lines[i-1].strip()):
                next_is_empty = (i == len(lines) - 1) or not lines[i+1].strip()
                if next_is_empty:
                    # 可能是标题
                    if line.endswith(':') or line.endswith('：'):
                        # 看起来像小节标题
                        improved_lines.append(f"### {line}")
                    else:
                        # 可能是普通标题
                        improved_lines.append(f"## {line}")
                    continue
            
            # 尝试处理列表项
            if line.startswith('•') or line.startswith('-') or (line[0].isdigit() and line[1:3] in ['. ', '、', ') ']):
                # 是列表项，直接添加
                improved_lines.append(line)
                continue
                
            # 普通文本行
            # 如果上一行不是空行，且不是以标点符号结束，可能是被分行的段落
            if current_paragraph and not current_paragraph[-1][-1] in ['.', '!', '?', '"', "'", '。', '！', '？', '"', '"', '…']:
                current_paragraph.append(line)
            else:
                # 开始新段落
                if current_paragraph:
                    improved_lines.append(' '.join(current_paragraph))
                current_paragraph = [line]
        
        # 处理最后一个段落
        if current_paragraph:
            improved_lines.append(' '.join(current_paragraph))
            
        return '\n\n'.join(improved_lines)
    
    def convert_to_markdown(self, content):
        """
        将PDF提取的文本内容转换为Markdown格式
        
        Args:
            content (str): PDF提取的文本内容
            
        Returns:
            str: Markdown格式的内容
        """
        # PDF提取的内容可能已经有了一些基本结构，需要进一步处理
        # 例如处理特殊字符，调整图片和表格等
        # 此处我们使用简单的标题和段落处理

        # 由于content已经在extract_content中被处理过，可以直接返回
        # 但也可以在这里进行进一步的Markdown标记添加
        
        # 特别处理一些常见问题
        # 替换连续多个换行为两个换行（Markdown段落）
        content = re.sub(r'\n{3,}', '\n\n', content)
        
        # 处理一些无法自动识别的列表
        # 例如将 "1. " 开头的行转换为Markdown列表
        content = re.sub(r'(?m)^(\d+)[\.\s]+\s*(.+)$', r'\1. \2', content)
        
        return content


class DocxConverter(FormatConverter):
    """DOCX格式转换器"""
    
    def extract_content(self, file_path):
        """
        从DOCX文件中提取内容
        
        Args:
            file_path (str): 文件路径
            
        Returns:
            tuple: (content, title) - 提取的内容和标题
        """
        if Document is None:
            print("警告: python-docx库未安装，无法提取DOCX内容")
            return "", os.path.splitext(os.path.basename(file_path))[0]
        
        try:
            doc = Document(file_path)
            title = os.path.splitext(os.path.basename(file_path))[0]
            
            # 提取所有段落的文本
            paragraphs = [para.text for para in doc.paragraphs]
            content = '\n'.join(paragraphs)
            
            # 尝试提取标题（通常是第一个非空段落）
            for para in doc.paragraphs:
                if para.text.strip():
                    title = para.text.strip()
                    break
            
            return content, title
        except Exception as e:
            print(f"提取DOCX内容出错: {e}")
            return "", os.path.splitext(os.path.basename(file_path))[0]
    
    def convert_to_markdown(self, content):
        """
        将DOCX提取的文本内容转换为Markdown格式
        
        Args:
            content (str): DOCX提取的文本内容
            
        Returns:
            str: Markdown格式的内容
        """
        # DOCX提取的内容已经是纯文本，需要进行一些格式化处理
        lines = content.split('\n')
        markdown_lines = []
        
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                markdown_lines.append("")
                continue
            
            # 尝试识别标题（第一行或短行）
            if i == 0 or (len(line) < 80 and (i == 0 or not lines[i-1].strip())):
                is_next_line_empty = (i == len(lines) - 1) or not lines[i+1].strip()
                
                if is_next_line_empty:
                    # 可能是标题
                    markdown_lines.append(f"## {line}")
                else:
                    markdown_lines.append(line)
            else:
                markdown_lines.append(line)
        
        return '\n'.join(markdown_lines)


class PPTXConverter(FormatConverter):
    """PPTX格式转换器"""
    
    def extract_content(self, file_path):
        """
        从PPTX文件中提取内容
        
        Args:
            file_path (str): 文件路径
            
        Returns:
            tuple: (content, title) - 提取的内容和标题
        """
        if Presentation is None:
            print("警告: python-pptx库未安装，无法提取PPTX内容")
            return "", os.path.splitext(os.path.basename(file_path))[0]
        
        try:
            prs = Presentation(file_path)
            title = os.path.splitext(os.path.basename(file_path))[0]
            
            # 提取所有幻灯片的文本
            all_text = []
            
            # 添加文档标题
            all_text.append(f"# {title}")
            all_text.append("")  # 空行
            
            for i, slide in enumerate(prs.slides):
                # 添加幻灯片标题
                slide_title = f"## 幻灯片 {i+1}"
                
                # 检查幻灯片是否有标题
                if slide.shapes.title and slide.shapes.title.text:
                    slide_title = f"## 幻灯片 {i+1}: {slide.shapes.title.text}"
                
                all_text.append(slide_title)
                all_text.append("")  # 空行
                
                # 提取幻灯片中的文本
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # 避免重复添加标题
                        if shape.text != slide.shapes.title.text if slide.shapes.title else True:
                            slide_text.append(shape.text)
                
                if slide_text:
                    all_text.extend(slide_text)
                    all_text.append("")  # 在幻灯片之间添加空行
            
            content = "\n".join(all_text)
            
            # 尝试提取更好的标题
            # 通常第一张幻灯片的标题可能是整个演示文稿的标题
            if prs.slides and prs.slides[0].shapes.title and prs.slides[0].shapes.title.text:
                title = prs.slides[0].shapes.title.text
            
            return content, title
            
        except Exception as e:
            print(f"提取PPTX内容出错: {e}")
            import traceback
            traceback.print_exc()
            return "", os.path.splitext(os.path.basename(file_path))[0]
    
    def convert_to_markdown(self, content):
        """
        将PPTX提取的文本内容转换为Markdown格式
        
        Args:
            content (str): PPTX提取的文本内容
            
        Returns:
            str: Markdown格式的内容
        """
        # PPTX内容在extract_content中已经组织为Markdown格式
        # 这里可以做一些额外的清理和格式优化
        
        # 移除多余的空行
        content = re.sub(r'\n{3,}', '\n\n', content)
        
        return content


class XLSXConverter(FormatConverter):
    """XLSX格式转换器"""
    
    def extract_content(self, file_path):
        """
        从XLSX文件中提取内容
        
        Args:
            file_path (str): 文件路径
            
        Returns:
            tuple: (content, title) - 提取的内容和标题
        """
        if openpyxl is None:
            print("警告: openpyxl库未安装，无法提取XLSX内容")
            return "", os.path.splitext(os.path.basename(file_path))[0]
        
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            title = os.path.splitext(os.path.basename(file_path))[0]
            
            # 提取所有工作表的内容
            all_text = []
            
            # 添加文档标题
            all_text.append(f"# {title}")
            all_text.append("")  # 空行
            
            for sheet in wb.worksheets:
                # 添加工作表标题
                all_text.append(f"## 工作表: {sheet.title}")
                all_text.append("")  # 空行
                
                # 提取表格内容并格式化为Markdown表格
                rows = []
                max_cols = 0
                
                # 计算表格的实际大小
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None and str(cell).strip() for cell in row):
                        rows.append(row)
                        max_cols = max(max_cols, len(row))
                
                if not rows:
                    all_text.append("*空工作表*")
                    all_text.append("")
                    continue
                
                # 创建Markdown表格
                # 添加表头行
                header = "| " + " | ".join([f"列 {i+1}" for i in range(max_cols)]) + " |"
                all_text.append(header)
                
                # 添加分隔行
                separator = "| " + " | ".join(["---" for _ in range(max_cols)]) + " |"
                all_text.append(separator)
                
                # 添加数据行
                for row in rows:
                    # 确保行有足够的列
                    padded_row = list(row) + ['' for _ in range(max_cols - len(row))]
                    # 将None转换为空字符串，并修正任何特殊字符
                    cleaned_row = [str(cell).replace('|', '\\|') if cell is not None else '' for cell in padded_row]
                    row_str = "| " + " | ".join(cleaned_row) + " |"
                    all_text.append(row_str)
                
                all_text.append("")  # 在工作表之间添加空行
            
            content = "\n".join(all_text)
            return content, title
            
        except Exception as e:
            print(f"提取XLSX内容出错: {e}")
            import traceback
            traceback.print_exc()
            return "", os.path.splitext(os.path.basename(file_path))[0]
    
    def convert_to_markdown(self, content):
        """
        将XLSX提取的文本内容转换为Markdown格式
        
        Args:
            content (str): XLSX提取的文本内容
            
        Returns:
            str: Markdown格式的内容
        """
        # XLSX内容在extract_content中已经组织为Markdown表格格式
        return content


class ConverterFactory:
    """转换器工厂类"""
    
    @staticmethod
    def get_converter(file_path):
        """
        根据文件路径获取对应的转换器
        
        Args:
            file_path (str): 文件路径
            
        Returns:
            FormatConverter: 对应格式的转换器实例
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.md', '.markdown']:
            return MarkdownConverter()
        elif ext in ['.html', '.htm']:
            return HtmlConverter()
        elif ext in ['.txt', '.text']:
            return TextConverter()
        elif ext == '.pdf':
            return PdfConverter()
        elif ext in ['.docx', '.doc']:
            return DocxConverter()
        elif ext in ['.pptx', '.ppt']:
            return PPTXConverter()
        elif ext in ['.xlsx', '.xls']:
            return XLSXConverter()
        else:
            # 默认使用文本转换器
            return TextConverter() 